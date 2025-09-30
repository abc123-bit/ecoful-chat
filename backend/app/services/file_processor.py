from __future__ import annotations

import io
import re
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Tuple, BinaryIO
from pathlib import Path
import zipfile
import tempfile

import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pandas as pd
from bs4 import BeautifulSoup
from loguru import logger
from app.core.config import settings

# ===== 本包依赖（统一 services 下）=====
# 可选配置：如无自定义 config，这里提供最小兜底
try:
    from .utils.config import get_config  # 确保 utils 是 services 的子包
except ImportError:  # 使用 ImportError 更具体
    def get_config() -> Dict[str, object]:
        return {
            "MINERU_ENABLED": True,
            "MINERU_POLL_SECONDS": 2,
            "MINERU_MAX_ATTEMPTS": 60,
            "MINERU_KEEP_MD": True,
        }

# MinerU 适配（可选存在）
try:
    from .mineru_parser import MinerUParser  # 你已移动到 services/
except Exception:
    MinerUParser = None  # type: ignore

try:
    from .mineru_batch_service import MinerUBatchService  # 若你使用批处理
except Exception:
    MinerUBatchService = None  # type: ignore
from .utils.config import get_config_manager
CFG = get_config_manager("./app/services/utils/config.yaml").get_all()

# ==================== 抽象基类 ====================
class FileProcessor(ABC):
    """文件处理器基类"""

    # 判断是否可以处理此类型文件
    @abstractmethod
    def can_process(self, file_type: str, mime_type: str) -> bool:
        """判断是否可以处理此类型文件"""
        pass

    # 提取文件文本内容
    @abstractmethod
    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        """
        提取文件文本内容

        Returns:
            Tuple[extracted_text, metadata]
        """
        pass

# ==================== MinerU 适配工具 ====================
def _extract_full_md(zf: zipfile.ZipFile) -> str:
    """在 ZIP 包中寻找 full.md（或等价）"""
    candidates = ['full.md', 'Full.md', 'FULL.MD', 'content.md', 'document.md']
    names = zf.namelist()
    lower_map = {n.lower(): n for n in names}
    for name in candidates:
        if name in lower_map:
            with zf.open(lower_map[name]) as fp:
                return fp.read().decode('utf-8', errors='ignore')
    # 兜底：取最大的 .md
    md_files = [n for n in names if n.lower().endswith('.md')]
    if md_files:
        md_files.sort(key=lambda n: zf.getinfo(n).file_size, reverse=True)
        with zf.open(md_files[0]) as fp:
            return fp.read().decode('utf-8', errors='ignore')
    return ""


def _read_full_md_from_zip_bytes(zip_bytes: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
        return _extract_full_md(zf)


def _read_full_md_from_zip_path(zip_path: Path) -> str:
    with zipfile.ZipFile(str(zip_path), 'r') as zf:
        return _extract_full_md(zf)


def _coerce_mineru_output(result, filename: str) -> Tuple[str, Dict[str, object]]:
    """把 MinerU 的多样返回统一为 (text, metadata)"""
    if result is None:
        return "", {}

    # 1) 直接文本
    if isinstance(result, str) and (('\n' in result) or len(result) > 0):
        return result, {"processor": "mineru", "source": "direct-text", "filename": filename}

    # 2) dict
    if isinstance(result, dict):
        meta = {k: v for k, v in result.items() if k not in {"zip_bytes", "zip_path", "full_md"}}
        if 'full_md' in result and isinstance(result['full_md'], str):
            return result['full_md'], {"processor": "mineru", "source": "dict-full_md", **meta}
        if 'zip_bytes' in result and isinstance(result['zip_bytes'], (bytes, bytearray)):
            text = _read_full_md_from_zip_bytes(result['zip_bytes'])
            return text, {"processor": "mineru", "source": "dict-zip_bytes", **meta}
        if 'zip_path' in result:
            text = _read_full_md_from_zip_path(Path(result['zip_path']))
            return text, {"processor": "mineru", "source": "dict-zip_path", **meta}

    # 3) (text, meta)
    if isinstance(result, tuple) and len(result) == 2 and isinstance(result[0], str):
        text, meta = result
        meta = meta or {}
        meta.setdefault("processor", "mineru")
        meta.setdefault("source", "tuple")
        meta.setdefault("filename", filename)
        return text, meta

    # 4) zip bytes
    if isinstance(result, (bytes, bytearray)):
        text = _read_full_md_from_zip_bytes(result)
        return text, {"processor": "mineru", "source": "zip-bytes", "filename": filename}

    # 5) zip path
    if isinstance(result, (str, Path)) and str(result).lower().endswith(".zip"):
        text = _read_full_md_from_zip_path(Path(result))
        return text, {"processor": "mineru", "source": "zip-path", "filename": filename}

    return "", {"processor": "mineru", "source": "unknown", "filename": filename}


def _stage_to_temp(file_stream: BinaryIO, suffix: str) -> Path:
    """把流写到临时文件，返回路径"""
    file_stream.seek(0)
    data = file_stream.read()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        tmp.write(data)
        tmp.flush()
        return Path(tmp.name)
    finally:
        tmp.close()


def _try_mineru(file_path: Path, filename: str) -> Tuple[str, Dict[str, object]]:
    """尝试用 MinerU 解析；失败返回空文本"""
    if not bool(CFG.get("MINERU_ENABLED", True)):
        return "", {}
    # 优先单文件解析器
    if MinerUParser is not None:
        try:
            parser = MinerUParser()  # 允许不同方法名
            for m in ("parse_document", "parse_file", "parse_single", "run", "process"):
                if hasattr(parser, m):
                    logger.info(f"MinerUParser.{m} -> {filename}")
                    result = getattr(parser, m)(str(file_path))
                    text, meta = _coerce_mineru_output(result, filename)
                    if text:
                        return text, meta
        except Exception as e:
            logger.warning(f"MinerUParser 失败（将尝试批处理回退）：{e}")
    # 批处理回退
    if MinerUBatchService is not None:
        try:
            svc = MinerUBatchService()
            for m in ("process_file", "process_single", "process_raw_files", "run"):
                if hasattr(svc, m):
                    logger.info(f"MinerUBatchService.{m} -> {filename}")
                    result = getattr(svc, m)(str(file_path))
                    text, meta = _coerce_mineru_output(result, filename)
                    if text:
                        return text, meta
        except Exception as e:
            logger.warning(f"MinerUBatchService 失败：{e}")

    return "", {}


# ==================== 具体处理器 ====================
# PDF文件处理器 
class PDFProcessor(FileProcessor):
    """PDF文件处理器"""
    """PDF文件处理器（优先 MinerU，失败本地回退）"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return mime_type == "application/pdf" or file_type.lower() == "pdf"

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        # 先尝试 MinerU
        try:
            tmp_path = _stage_to_temp(file_stream, ".pdf")
            mineru_text, mineru_meta = _try_mineru(tmp_path, filename)
            if mineru_text:
                cleaned_len = len(mineru_text)
                mineru_meta.update({
                    "text_length": cleaned_len,
                    "extracted_pages": None,  # MinerU 没有页级计数时置 None
                })
                return mineru_text, mineru_meta
        except Exception as e:
            logger.warning(f"MinerU 解析 PDF 失败，将用本地回退：{e}")

        # 回退：沿用你原有本地提取（保持原逻辑/字段）

        try:
            pdf_reader = PyPDF2.PdfReader(file_stream)
            text_content = []
            metadata = {
                "total_pages": len(pdf_reader.pages),
                "title": None,
                "author": None,
                "creator": None,
                "subject": None
            }

            # 提取元数据
            if pdf_reader.metadata:
                metadata.update({
                    "title": pdf_reader.metadata.get("/Title"),
                    "author": pdf_reader.metadata.get("/Author"),
                    "creator": pdf_reader.metadata.get("/Creator"),
                    "subject": pdf_reader.metadata.get("/Subject")
                })

            # 提取每页文本
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"# 第{page_num + 1}页\n\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")

            extracted_text = "\n\n".join(text_content)
            metadata["extracted_pages"] = len(text_content)
            metadata["text_length"] = len(extracted_text)

            return extracted_text, metadata

        except Exception as e:
            logger.error(f"Failed to process PDF file {filename}: {e}")
            raise

# word文档处理器(docx\doc)
class DocxProcessor(FileProcessor):
    """Word文档处理器"""
    """Word文档处理器（优先 MinerU，失败本地回退）"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return (
            mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or file_type.lower() in ["docx", "doc"]
        )

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
         # 先尝试 MinerU
        try:
            suffix = ".docx" if filename.lower().endswith(".docx") else ".doc"
            tmp_path = _stage_to_temp(file_stream, suffix)
            mineru_text, mineru_meta = _try_mineru(tmp_path, filename)
            if mineru_text:
                mineru_meta.update({
                    "text_length": len(mineru_text),
                })
                return mineru_text, mineru_meta
        except Exception as e:
            logger.warning(f"MinerU 解析 Word 失败，将用本地回退：{e}")

        # 回退：沿用你原有本地提取（保持原逻辑/字段）
        
        try:
            # 注意：python-docx 不支持 .doc；如为 .doc 这里可能抛错，交给上层处理
            doc = Document(file_stream)
            text_content = []
            metadata = {
                "total_paragraphs": len(doc.paragraphs),
                "total_tables": len(doc.tables),
                "title": None,
                "author": None,
                "subject": None
            }

            # 提取核心属性
            if doc.core_properties:
                metadata.update({
                    "title": doc.core_properties.title,
                    "author": doc.core_properties.author,
                    "subject": doc.core_properties.subject
                })

            # 提取段落文本
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text.strip())

            # 提取表格内容
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    text_content.append(f"\\n\\n**表格内容:**\\n{table_text}")

            extracted_text = "\\n\\n".join(text_content)
            metadata["text_length"] = len(extracted_text)

            return extracted_text, metadata

        except Exception as e:
            logger.error(f"Failed to process DOCX file {filename}: {e}")
            raise

    def _extract_table_text(self, table) -> str:
        """提取表格文本"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(" | ".join(row_data))
        return "\\n".join(table_data)

# Excel文件处理器("xlsx", "xls")
class ExcelProcessor(FileProcessor):
    """Excel文件处理器"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return (
            mime_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel"
            ]
            or file_type.lower() in ["xlsx", "xls"]
        )

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        try:
            # 使用pandas读取Excel文件
            workbook = pd.read_excel(file_stream, sheet_name=None, engine='openpyxl')
            text_content = []
            metadata = {
                "total_sheets": len(workbook),
                "sheet_names": list(workbook.keys())
            }

            for sheet_name, df in workbook.items():
                if not df.empty:
                    sheet_text = f"**工作表: {sheet_name}**\\n\\n"

                    # 转换为文本表格
                    sheet_text += df.to_string(index=False, na_rep="")
                    text_content.append(sheet_text)

            extracted_text = "\\n\\n".join(text_content)
            metadata["text_length"] = len(extracted_text)

            return extracted_text, metadata

        except Exception as e:
            logger.error(f"Failed to process Excel file {filename}: {e}")
            raise

# PowerPoint文件处理器
class PowerPointProcessor(FileProcessor):
    """PowerPoint文件处理器"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return (
            mime_type in [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint"
            ]
            or file_type.lower() in ["pptx", "ppt"]
        )

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        try:
            prs = Presentation(file_stream)
            text_content = []
            metadata = {
                "total_slides": len(prs.slides),
                "slide_titles": []
            }

            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"# 幻灯片 {slide_num + 1}\\n\\n"
                slide_content = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                if slide_content:
                    slide_text += "\\n\\n".join(slide_content)
                    text_content.append(slide_text)

                    # 尝试提取标题
                    if slide_content:
                        metadata["slide_titles"].append(slide_content[0][:50])

            extracted_text = "\\n\\n".join(text_content)
            metadata["text_length"] = len(extracted_text)

            return extracted_text, metadata

        except Exception as e:
            logger.error(f"Failed to process PowerPoint file {filename}: {e}")
            raise

# 纯文本文件处理器("txt", "md", "markdown", "csv")
class TextProcessor(FileProcessor):
    """纯文本文件处理器"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return (
            mime_type.startswith("text/")
            or file_type.lower() in ["txt", "md", "markdown", "csv"]
        )

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        try:
            # 尝试不同编码
            encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
            content = None

            for encoding in encodings:
                try:
                    file_stream.seek(0)
                    content = file_stream.read().decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                raise ValueError("Cannot decode file with any supported encoding")

            metadata = {
                "file_type": Path(filename).suffix.lower(),
                "text_length": len(content),
                "line_count": len(content.splitlines())
            }

            # 如果是CSV文件，尝试格式化
            if filename.lower().endswith('.csv'):
                try:
                    file_stream.seek(0)
                    df = pd.read_csv(file_stream)
                    content = df.to_string(index=False)
                    metadata["csv_rows"] = len(df)
                    metadata["csv_columns"] = len(df.columns)
                except:
                    pass  # 如果CSV解析失败，使用原始文本

            return content, metadata

        except Exception as e:
            logger.error(f"Failed to process text file {filename}: {e}")
            raise

# 图片文件处理器(仅提取基本信息，未来可扩展OCR)
class ImageProcessor(FileProcessor):
    """图片文件处理器"""

    def can_process(self, file_type: str, mime_type: str) -> bool:
        return mime_type.startswith("image/")

    def extract_text(self, file_stream: BinaryIO, filename: str) -> Tuple[str, Dict]:
        try:
            image = Image.open(file_stream)
            metadata = {
                "format": image.format,
                "mode": image.mode,
                "size": image.size,
                "width": image.width,
                "height": image.height
            }

            # 对于图片，返回基本信息作为"文本"
            extracted_text = (
                f"图片文件: {filename}\\n"
                f"格式: {image.format}\\n"
                f"尺寸: {image.width} x {image.height}\\n"
                f"颜色模式: {image.mode}"
            )

            # TODO: 如果需要OCR功能，可以在这里集成OCR库
            # 例如使用 pytesseract 进行文字识别

            metadata["text_length"] = len(extracted_text)
            return extracted_text, metadata

        except Exception as e:
            logger.error(f"Failed to process image file {filename}: {e}")
            raise

# 统一文件处理器管理器，根据不同类型调用不同处理器
class FileProcessorManager:
    """文件处理器管理器"""

    # 初始化时注册所有处理器
    def __init__(self):
        self.processors = [
            PDFProcessor(),
            DocxProcessor(),
            ExcelProcessor(),
            PowerPointProcessor(),
            TextProcessor(),
            ImageProcessor()
        ]

    # 根据文件类型获取合适的处理器
    def get_processor(self, file_type: str, mime_type: str) -> Optional[FileProcessor]:
        """根据文件类型获取合适的处理器"""
        for processor in self.processors:
            if processor.can_process(file_type, mime_type):
                return processor
        return None

    # 处理文件并提取文本
    def process_file(
        self,
        file_stream: BinaryIO,
        filename: str,
        mime_type: str
    ) -> Tuple[str, Dict]:
        """
        处理文件并提取文本

        Args:
            file_stream: 文件流
            filename: 文件名
            mime_type: MIME类型

        Returns:
            Tuple[extracted_text, metadata]

        Raises:
            ValueError: 不支持的文件类型
            Exception: 处理过程中的错误
        """
        file_type = Path(filename).suffix[1:].lower()  # 去掉点号
        processor = self.get_processor(file_type, mime_type)

        if processor is None:
            raise ValueError(f"Unsupported file type: {file_type} ({mime_type})")

        try:
            file_stream.seek(0)  # 确保从文件开头读取
            extracted_text, metadata = processor.extract_text(file_stream, filename)

            # 清理文本
            cleaned_text = self._clean_text(extracted_text)

            # 更新元数据
            metadata.update({
                "processor_type": processor.__class__.__name__,
                "original_filename": filename,
                "mime_type": mime_type,
                "file_type": file_type,
                "cleaned_text_length": len(cleaned_text)
            })

            return cleaned_text, metadata

        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            raise

    # 清理提取的文本
    def _clean_text(self, text: str) -> str:
        """清理提取的文本"""
        if not text:
            return ""

        # 移除多余的空白字符
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # 多个连续换行替换为两个
        text = re.sub(r'[ \t]+', ' ', text)  # 多个空格替换为一个
        text = re.sub(r'^\s+|\s+$', '', text, flags=re.MULTILINE)  # 移除行首行尾空格
        text = re.sub(r'\u3000', ' ', text)  # 全角空格替换为半角空格
        text = re.sub(r'\xa0', ' ', text)  # 不间断空格替换为普通空格
        
        # 移除特殊字符
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
        text = re.sub(r'[\x7f-\x9f]', '', text)

        return text.strip()

    # 获取支持的文件类型
    def get_supported_types(self) -> Dict[str, List[str]]:
        """获取支持的文件类型"""
        return {
            "PDF": ["pdf"],
            "Word文档": ["doc", "docx"],
            "Excel表格": ["xls", "xlsx"],
            "PowerPoint演示文稿": ["ppt", "pptx"],
            "纯文本": ["txt", "md", "markdown", "csv"],
            "图片": ["jpg", "jpeg", "png", "gif", "bmp", "tiff"]
        }

# 创建全局文件处理器实例
file_processor = FileProcessorManager()